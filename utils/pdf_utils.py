import pdfplumber
from pptx import Presentation
from pathlib import Path
import re
import os
from sentence_transformers import SentenceTransformer

# Matches the "lOMoARcPSD|<id>" line StuDocu injects into the text layer of
# downloaded documents as a tracking watermark — garbage that otherwise ends
# up as slide/question content (e.g. becoming a topic title verbatim).
_WATERMARK_LINE_RE = re.compile(r"(?im)^\s*lOMoARcPSD\s*\|\s*\d+\s*$")

def extract_text_from_pdf(pdf_path: str) -> list[dict]:
    pages = []
    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text:
                text = _WATERMARK_LINE_RE.sub("", text).strip()
            if text and text.strip():
                pages.append({"page": i + 1, "text": text.strip()})
    return pages

def extract_text_from_pptx(pptx_path: str) -> list[dict]:
    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            slides.append({"page": i + 1, "text": "\n".join(texts)})
    return slides

def chunk_slides(pages: list[dict], chunk_size: int = 3,
                 overlap: int = 0) -> list[dict]:
    """Group consecutive slide pages into chunks.

    ``chunk_size`` pages per chunk; ``overlap`` pages are shared with the
    previous chunk (a sliding page window). ``overlap=0`` (the default) is the
    original non-overlapping behaviour and yields identical chunk ids
    (``slides_0``, ``slides_1``, …) to callers that don't set it.
    """
    if chunk_size < 1:
        raise ValueError(f"chunk_size must be >= 1, got {chunk_size}")
    if not 0 <= overlap < chunk_size:
        raise ValueError(
            f"overlap must satisfy 0 <= overlap < chunk_size, "
            f"got overlap={overlap}, chunk_size={chunk_size}"
        )
    step = chunk_size - overlap
    chunks = []
    cid = 0
    for i in range(0, len(pages), step):
        group = pages[i:i + chunk_size]
        if not group:
            break
        combined_text = "\n\n".join(p["text"] for p in group)
        page_nums = [p["page"] for p in group]
        chunks.append({
            "chunk_id": f"slides_{cid}",
            "text": combined_text,
            "pages": page_nums,
            "source": "slides"
        })
        cid += 1
        # With overlap, the window can reach the end before the range does;
        # stop so we don't emit a trailing chunk fully contained in the last.
        if i + chunk_size >= len(pages):
            break
    return chunks


def chunk_slides_by_tokens(pages: list[dict], tokenizer, max_tokens: int = 512,
                           overlap: int = 64) -> list[dict]:
    """Token-window chunker: concatenate page text and split every
    ``max_tokens`` tokens with ``overlap`` tokens shared between neighbours.

    ``tokenizer`` is a HuggingFace-style tokenizer (e.g. the embedding model's
    ``.tokenizer``) so token boundaries match how the text will be embedded.
    Each chunk records the source ``pages`` its tokens came from, so gold
    labels anchored to page numbers can be re-mapped across configs.
    """
    if max_tokens < 1:
        raise ValueError(f"max_tokens must be >= 1, got {max_tokens}")
    if not 0 <= overlap < max_tokens:
        raise ValueError(
            f"overlap must satisfy 0 <= overlap < max_tokens, "
            f"got overlap={overlap}, max_tokens={max_tokens}"
        )
    flat_ids: list[int] = []
    token_page: list[int] = []  # source page number for each token position
    for p in pages:
        ids = tokenizer.encode(p["text"], add_special_tokens=False)
        flat_ids.extend(ids)
        token_page.extend([p["page"]] * len(ids))

    step = max_tokens - overlap
    chunks = []
    cid = 0
    n = len(flat_ids)
    for start in range(0, n, step):
        end = min(start + max_tokens, n)
        span = flat_ids[start:end]
        if not span:
            break
        chunks.append({
            "chunk_id": f"slides_{cid}",
            "text": tokenizer.decode(span),
            "pages": sorted(set(token_page[start:end])),
            "source": "slides"
        })
        cid += 1
        if end >= n:
            break
    return chunks

def chunk_past_paper(pages: list[dict], paper_year: str,
                     paper_id: str | None = None) -> list[dict]:
    # ``paper_id`` namespaces the chunk ids so two papers sharing a detected
    # year (e.g. sem1 + sem2 of 2022) don't collide on ``paper_2022_q1``. It
    # defaults to the year, keeping the id format unchanged for the common
    # one-paper-per-year case. The stored ``year`` is always the real year.
    id_ns = paper_id if paper_id is not None else paper_year
    full_text = "\n".join(p["text"] for p in pages)
    question_pattern = re.compile(
        r'(?=^(?:Q\s*\d+|Question\s+\d+|\d+[\.\)]\s))',
        re.MULTILINE | re.IGNORECASE
    )
    parts = question_pattern.split(full_text)
    parts = [p.strip() for p in parts if len(p.strip()) > 50]
    chunks = []
    for i, part in enumerate(parts):
        chunks.append({
            "chunk_id": f"paper_{id_ns}_q{i+1}",
            "text": part,
            "year": paper_year,
            "question_num": i + 1,
            "source": "past_paper"
        })
    return chunks


# ── Test-facing parse / flag helpers ────────────────────────

# Lazy-loaded embedder for duplicate detection
_embedder = None


def _get_embedder():
    global _embedder
    if _embedder is None:
        _embedder = SentenceTransformer("all-MiniLM-L6-v2")
    return _embedder


def parse_pdf(file_path: str, type: str = "past_paper") -> list[dict]:
    """Parse a PDF into question-chunk dicts.

    ``file_path`` is relative to the project root (e.g. ``fixtures/paper_2022.pdf``).
    For past papers the file is split on question boundaries; each returned dict
    has keys ``chunk_id``, ``text``, ``year``, ``question_num``, ``source``,
    and ``is_repeat`` (default ``False``).
    """
    # Resolve path relative to project root
    project_root = Path(__file__).resolve().parent.parent
    abs_path = project_root / file_path

    pages = extract_text_from_pdf(str(abs_path))

    # Derive year — try PDF text first, then filename
    full_text = "\n".join(p["text"] for p in pages)
    year = _extract_year_from_text(full_text)
    if year is None:
        year = _extract_year_from_filename(str(abs_path))
    if year is None:
        year = "unknown"
    else:
        year = str(year)

    if type == "past_paper":
        chunks = chunk_past_paper(pages, year)
    elif type == "slides":
        chunks = chunk_slides(pages)
    else:
        chunks = chunk_past_paper(pages, year)

    for c in chunks:
        c.setdefault("is_repeat", False)
    return chunks


def flag_duplicate_questions(questions: list[dict]) -> list[dict]:
    """Mark cross-paper duplicate questions with ``is_repeat = True``.

    Returns the same list (mutated in-place AND returned).  A question is
    flagged as a repeat when it is semantically similar to a question from a
    different paper year.  Similarity is measured via cosine distance on
    ``all-MiniLM-L6-v2`` embeddings (threshold < 0.6).
    """
    if len(questions) < 2:
        return questions

    embedder = _get_embedder()
    texts = [q["text"][:2000] for q in questions]
    years = [q.get("year", "") for q in questions]
    embeddings = embedder.encode(texts, show_progress_bar=False)

    SIMILARITY_THRESHOLD = 0.3  # cosine similarity (1 - cosine distance)

    for i in range(len(questions)):
        for j in range(i + 1, len(questions)):
            if years[i] == years[j] or not years[i] or not years[j]:
                continue
            # Cosine similarity = dot product (embeddings are L2-normalized)
            sim = float(embeddings[i] @ embeddings[j])
            if sim >= SIMILARITY_THRESHOLD:
                # Mark the one from the later year as repeat
                if years[i] < years[j]:
                    questions[j]["is_repeat"] = True
                else:
                    questions[i]["is_repeat"] = True

    return questions


def _extract_year_from_text(text: str) -> int | None:
    """Try to extract an academic year ending from PDF body text.

    Handles patterns like ``Academic Year 2022/2023`` (→ 2023) and
    ``AY2019/2020`` (→ 2020).
    """
    # "Academic Year 2022/2023"
    m = re.search(r"Academic\s+Year\s+(\d{4})\s*/\s*(\d{4})", text, re.IGNORECASE)
    if m:
        return int(m.group(2))
    # "AY2019/2020"
    m = re.search(r"AY\s*(\d{4})\s*/\s*(\d{4})", text, re.IGNORECASE)
    if m:
        return int(m.group(2))
    return None


def _extract_year_from_filename(filepath: str) -> int | None:
    """Extract a 4-digit year (20xx) from a filename stem."""
    stem = os.path.splitext(os.path.basename(filepath))[0]
    m = re.search(r"(20\d{2})", stem)
    if m:
        return int(m.group(1))
    return None


def extract_format(file_paths: list[str]) -> dict:
    """Examine PDFs to determine the exam source year and how it was found.

    Returns ``{"source_year": int, "date_source": "pdf_text"|"filename"}``.

    Priority:
    1. Scan each PDF body for an academic-year pattern.
    2. Fall back to the first filename that contains a ``20xx`` year.
    """
    project_root = Path(__file__).resolve().parent.parent

    for fp in file_paths:
        abs_path = project_root / fp
        pages = extract_text_from_pdf(str(abs_path))
        full_text = "\n".join(p["text"] for p in pages)
        year = _extract_year_from_text(full_text)
        if year is not None:
            return {"source_year": year, "date_source": "pdf_text"}

    # Fallback to filename
    for fp in file_paths:
        year = _extract_year_from_filename(str(fp))
        if year is not None:
            return {"source_year": year, "date_source": "filename"}

    return {"source_year": None, "date_source": "none"}